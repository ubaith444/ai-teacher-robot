"""
core/document_loader.py
========================
Zoro Robot — Unified Document Loader
Loads and chunks documents from all supported sources.
Merged from v1 DocumentLoader + DocumentChunker.

Supports: PDF, TXT, Markdown, JSON (Q&A), CSV
Auto-detects: subject, source_type, grade from filename
"""

from __future__ import annotations

import hashlib
import json
import logging
import re
from pathlib import Path
from typing import Any

from app.schemas.contracts import DocumentChunk, SourceType

logger = logging.getLogger("zoro.loader")


# ═════════════════════════════════════════════════════════════
# DOCUMENT LOADER
# ═════════════════════════════════════════════════════════════


class DocumentLoader:
    """
    Loads documents from file or directory into raw {text, metadata} dicts.
    Handles PDF (PyMuPDF preferred, pdfplumber fallback), TXT, MD, JSON, CSV.
    """

    SUPPORTED = {".pdf", ".txt", ".md", ".json", ".csv", ".docx", ".pptx", ".doc", ".ppt"}

    SUBJECT_KEYWORDS = {
        "math": "Math",
        "algebra": "Math",
        "geometry": "Math",
        "calculus": "Math",
        "science": "Science",
        "physics": "Physics",
        "chemistry": "Chemistry",
        "biology": "Biology",
        "history": "History",
        "geography": "Geography",
        "english": "English",
        "hindi": "Hindi",
        "social": "Social",
        "ncert": "General",
    }

    GRADE_PATTERN = re.compile(r"class[_\s-]?(\d{1,2})", re.IGNORECASE)

    def load(self, source_path: str) -> list[dict[str, Any]]:
        path = Path(source_path)
        if path.is_dir():
            return self._load_directory(path)
        return self._load_file(path)

    def _load_directory(self, directory: Path) -> list[dict[str, Any]]:
        docs = []
        for p in sorted(directory.rglob("*")):
            if p.suffix.lower() in self.SUPPORTED:
                try:
                    loaded = self._load_file(p)
                    docs.extend(loaded)
                    logger.info(f"Loaded {len(loaded)} page(s) from {p.name}")
                except Exception as e:
                    logger.warning(f"Skipped {p.name}: {e}")
        logger.info(f"Total documents loaded: {len(docs)}")
        return docs

    def _load_file(self, path: Path) -> list[dict[str, Any]]:
        ext = path.suffix.lower()
        meta = self._build_meta(path)

        if ext == ".pdf":
            return self._load_pdf(path, meta)
        elif ext in (".txt", ".md"):
            text = path.read_text(encoding="utf-8", errors="ignore").strip()
            return [{"text": text, **meta}] if text else []
        elif ext == ".json":
            return self._load_json_qa(path, meta)
        elif ext == ".csv":
            return self._load_csv(path, meta)
        elif ext == ".docx":
            return self._load_docx(path, meta)
        elif ext == ".pptx":
            return self._load_pptx(path, meta)
        return []

    def _load_pdf(self, path: Path, meta: dict) -> list[dict]:
        pages = []
        try:
            import fitz

            doc = fitz.open(str(path))
            for i, page in enumerate(doc):
                text = page.get_text("text").strip()
                if text:
                    pages.append({"text": text, "page": i + 1, **meta})
            doc.close()
            return pages
        except ImportError:
            pass
        try:
            import pdfplumber

            with pdfplumber.open(str(path)) as pdf:
                for i, page in enumerate(pdf.pages):
                    text = (page.extract_text() or "").strip()
                    if text:
                        pages.append({"text": text, "page": i + 1, **meta})
            return pages
        except ImportError:
            logger.warning("No PDF library found. Install pymupdf or pdfplumber.")
        return pages

    def _load_docx(self, path: Path, meta: dict) -> list[dict]:
        try:
            import docx
            doc = docx.Document(str(path))
            full_text = []
            for para in doc.paragraphs:
                if para.text.strip():
                    full_text.append(para.text.strip())
            text = "\n".join(full_text)
            return [{"text": text, **meta}] if text else []
        except Exception as e:
            logger.warning(f"Error loading DOCX {path}: {e}")
            return []

    def _load_pptx(self, path: Path, meta: dict) -> list[dict]:
        try:
            from pptx import Presentation
            prs = Presentation(str(path))
            slides = []
            for i, slide in enumerate(prs.slides):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                text = "\n".join(slide_text)
                if text:
                    slides.append({"text": text, "page": i + 1, **meta})
            return slides
        except Exception as e:
            logger.warning(f"Error loading PPTX {path}: {e}")
            return []

    def _load_json_qa(self, path: Path, meta: dict) -> list[dict]:
        data = json.loads(path.read_text(encoding="utf-8"))
        docs = []
        if isinstance(data, list):
            for item in data:
                q = item.get("question", "")
                a = item.get("answer", "")
                if q and a:
                    docs.append({"text": f"Q: {q}\nA: {a}", **meta})
        elif isinstance(data, dict):
            # Support {question: answer} flat maps
            for q, a in data.items():
                docs.append({"text": f"Q: {q}\nA: {a}", **meta})
        return docs

    def _load_csv(self, path: Path, meta: dict) -> list[dict]:
        try:
            import pandas as pd

            df = pd.read_csv(path)
            docs = []
            if "text" in df.columns:
                for text in df["text"].dropna():
                    docs.append({"text": str(text), **meta})
            elif "question" in df.columns and "answer" in df.columns:
                for _, row in df.iterrows():
                    docs.append(
                        {"text": f"Q: {row['question']}\nA: {row['answer']}", **meta}
                    )
            return docs
        except ImportError:
            logger.warning("pandas not installed; skipping CSV.")
        return []

    def _build_meta(self, path: Path) -> dict:
        name = path.name.lower()
        subject = "General"
        for kw, subj in self.SUBJECT_KEYWORDS.items():
            if kw in name:
                subject = subj
                break

        grade_match = self.GRADE_PATTERN.search(name)
        grade = int(grade_match.group(1)) if grade_match else 0

        if "ncert" in name or "textbook" in name:
            src_type = SourceType.TEXTBOOK.value
        elif "notes" in name:
            src_type = SourceType.NOTES.value
        elif path.suffix == ".json" or "qa" in name:
            src_type = SourceType.QA.value
        elif path.suffix == ".pdf":
            src_type = SourceType.PDF.value
        elif path.suffix == ".csv":
            src_type = SourceType.CSV.value
        else:
            src_type = SourceType.DOCUMENT.value

        return {
            "source": path.name,
            "source_path": str(path),
            "source_type": src_type,
            "subject": subject,
            "grade": grade,
            "topic": "",
            "page": 0,
        }


# ═════════════════════════════════════════════════════════════
# DOCUMENT CHUNKER
# ═════════════════════════════════════════════════════════════


class DocumentChunker:
    """
    Splits raw document text into overlapping word-bounded chunks.
    Target: 80–280 words per chunk, 30-word sentence-boundary overlap.
    Produces DocumentChunk objects ready for embedding.
    """

    def __init__(
        self,
        min_words: int = 80,
        max_words: int = 280,
        overlap_words: int = 30,
    ):
        self.min_words = min_words
        self.max_words = max_words
        self.overlap_words = overlap_words

    def chunk(self, docs: list[dict[str, Any]]) -> list[DocumentChunk]:
        all_chunks: list[DocumentChunk] = []
        for doc in docs:
            text = doc.get("text", "").strip()
            if not text:
                continue
            sentences = self._split_sentences(text)
            groups = self._group_sentences(sentences)
            for idx, chunk_text in enumerate(groups):
                cid = self._make_id(doc.get("source", ""), idx)
                all_chunks.append(
                    DocumentChunk(
                        chunk_id=cid,
                        text=chunk_text,
                        source=doc.get("source", "unknown"),
                        source_type=doc.get("source_type", SourceType.DOCUMENT.value),
                        subject=doc.get("subject", ""),
                        topic=doc.get("topic", ""),
                        grade=doc.get("grade", 0),
                        page=doc.get("page", 0),
                        word_count=len(chunk_text.split()),
                    )
                )
        logger.info(f"Chunked {len(docs)} docs → {len(all_chunks)} chunks")
        return all_chunks

    def _split_sentences(self, text: str) -> list[str]:
        text = re.sub(r"\s+", " ", text)
        return [s.strip() for s in re.split(r"(?<=[.!?])\s+", text) if s.strip()]

    def _group_sentences(self, sentences: list[str]) -> list[str]:
        chunks, current, overlap_buf = [], [], []
        for sent in sentences:
            words = sent.split()
            if (
                len(current) + len(words) > self.max_words
                and len(current) >= self.min_words
            ):
                chunks.append(" ".join(current))
                current = overlap_buf[-self.overlap_words :] + words
                overlap_buf = words
            else:
                current.extend(words)
                overlap_buf = words
        if current and len(current) >= self.min_words // 2:
            chunks.append(" ".join(current))
        return chunks

    @staticmethod
    def _make_id(source: str, idx: int) -> str:
        return hashlib.md5(f"{source}::{idx}".encode()).hexdigest()[:12]

